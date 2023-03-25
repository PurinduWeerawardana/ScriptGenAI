from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import os



slide_count, image_count, table_count = 0, 0, 0
filename = 'Xpath.pptx'

def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape
                
def get_tables_from_presentation(pres):
   tables = list()
   for slide in pres.slides:
      for shp in iter(slide.shapes):
         
         if shp.has_table:
            table = shp.table
            tables.append(table)
            
   return tables


def iter_to_nonempty_table_cells(tbl):
   for ridx in range(sum(1 for _ in iter(tbl.rows))):
      for cidx in range(sum(1 for _ in iter(tbl.columns))):
         cell = tbl.cell(ridx, cidx)
         txt = type("")(cell.text)
         txt = txt.strip()
         yield txt


for picture in iter_picture_shapes(Presentation(filename)):
    image = picture.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    image_filename = 'image'+str(slide_count)+'.%s' % image.ext
    image_count = image_count+1
    slide_count = slide_count +1
        
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)
        #print("Image Saved", image_count, slide_count)
        
slide_count = 0               
for eachfile in glob.glob(filename):            
    prs = Presentation(eachfile)
    for slide in prs.slides:
        print("----------------------")
        print("S: ", slide_count)
        slide_count = slide_count+1


        #Find the heading paragraph in the slide
        heading=None
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.startswith("#"):
                    heading=shape.text.strip("#")
                    break
        if heading is not None:
            print("Heading:",heading)

        #find the sub points for the heading

            sub_points=[]
            for shape in slide.shapes:
                if hasattr(shape,"text"):
                    if shape.text.startswith("-"):
                        sub_points.append(shape.textstrip("-"))

            if len(sub_points)>0:
                print("Sub Points:")
                for sub_point in sub_points:
                    print("- " +sub_point)

        #print the text in shapes
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                if not shape.text.startswith("#")and not shape.text.startswith("-"):
                    print(shape.text)
                    print()

       
# extract tables from slide-show presentation
tables = get_tables_from_presentation(prs)

for tbl in tables:
   table_count = table_count+1
   print(len(tbl.columns))
   it = iter_to_nonempty_table_cells(tbl)
   for i in range(len(tbl.columns)):
      print("|".join(it))

print("There are \n Slides: ", slide_count, "\n Slides that have images: ",image_count, "\n Slides that have tables", table_count)



