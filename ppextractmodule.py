from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob


def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape


def get_tables_from_presentation(pres):
    tables = list()
    for slide in pres.slides:
        for shp in iter(slide.shapes):

            if shp.has_table:
                table = shp.table
                tables.append(table)

    return tables


def iter_to_nonempty_table_cells(tbl):
    for ridx in range(sum(1 for _ in iter(tbl.rows))):
        for cidx in range(sum(1 for _ in iter(tbl.columns))):
            cell = tbl.cell(ridx, cidx)
            txt = type("")(cell.text)
            txt = txt.strip()
            yield txt


def process(filename):
    slide_count, image_count, table_count = 0, 0, 0
    out = ""
    for picture in iter_picture_shapes(Presentation(filename)):
        image = picture.image
        # ---get image "file" contents---
        image_bytes = image.blob
        # ---make up a name for the file, e.g. 'image.jpg'---
        image_filename = 'image'+str(slide_count)+'.%s' % image.ext
        image_count = image_count+1
        slide_count = slide_count + 1

        with open(image_filename, 'wb') as f:
            f.write(image_bytes)
            # print("Image Saved", image_count, slide_count)

    slide_count = 0
    for eachfile in glob.glob(filename):
        prs = Presentation(eachfile)
        for slide in prs.slides:
            print("----------------------")
            out = out + "----------------------" + "\n"
            print("S: ", slide_count)
            out = out + "S:" + str(slide_count) + "\n"
            slide_count = slide_count+1

            # Find the heading paragraph in the slide
            heading = None
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.startswith("#"):
                        heading = shape.text.strip("#")
                        break
            if heading is not None:
                print("Heading:", heading)
                out = out + "Heading: " + heading + "\n"

            # find the sub points for the heading

                sub_points = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.startswith("-"):
                            sub_points.append(shape.textstrip("-"))

                if len(sub_points) > 0:
                    print("Sub Points:")
                    out = out + "Sub Points: " + "\n"
                    for sub_point in sub_points:
                        print("- " + sub_point)
                        out = out + "- " + sub_point + "\n"

            # print the text in shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if not shape.text.startswith("#") and not shape.text.startswith("-"):
                        print(shape.text)
                        out = out + shape.text + "\n\n"
                        print()
    # extract tables from slide-show presentation
    tables = get_tables_from_presentation(prs)
    for tbl in tables:
        table_count = table_count+1
        print(len(tbl.columns))
        out = out + str(len(tbl.columns)) + "\n"
        it = iter_to_nonempty_table_cells(tbl)
        for i in range(len(tbl.columns)):
            print("|".join(it))
            out = out + "|".join(it) + "\n"

    out = out + "There are \n Slides: " + str(slide_count) + "\n Slides that have images: " + str(
    image_count) + "\n Slides that have tables" + str(table_count) + ".\n"
    return {"text" : out, "images": image_count}